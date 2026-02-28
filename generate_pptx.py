"""
Generate ToledoMessage PowerPoint Presentation
Rich diagrams, animations, graphs, and professional UI/UX
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from io import BytesIO

# ─── Color Palette ───
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_NAVY     = RGBColor(0x16, 0x21, 0x3E)
C_BLUE     = RGBColor(0x19, 0x76, 0xD2)
C_BLUE_L   = RGBColor(0x42, 0xA5, 0xF5)
C_TEAL     = RGBColor(0x00, 0x80, 0x69)
C_GREEN    = RGBColor(0x25, 0xD3, 0x66)
C_ORANGE   = RGBColor(0xFF, 0x98, 0x00)
C_RED      = RGBColor(0xE5, 0x39, 0x35)
C_PURPLE   = RGBColor(0x7C, 0x4D, 0xFF)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY_L   = RGBColor(0xE0, 0xE0, 0xE0)
C_GRAY     = RGBColor(0x9E, 0x9E, 0x9E)
C_GRAY_D   = RGBColor(0x42, 0x42, 0x42)
C_BG       = RGBColor(0x0F, 0x11, 0x1A)
C_ACCENT   = RGBColor(0x3A, 0x76, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Helper Functions ───

def add_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, c1=C_BG, c2=C_DARK):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0

def add_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=C_WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=C_WHITE,
                   line_spacing=1.3, font_name='Segoe UI', bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz if sz else font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bld if bld else False
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_arrow_line(slide, x1, y1, x2, y2, color=C_BLUE_L, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=C_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def make_chart_image(chart_func, w=6, h=4, dpi=150):
    buf = BytesIO()
    fig, ax = plt.subplots(figsize=(w, h), dpi=dpi)
    fig.patch.set_facecolor('#0F111A')
    ax.set_facecolor('#0F111A')
    chart_func(fig, ax)
    fig.savefig(buf, format='png', bbox_inches='tight', facecolor='#0F111A', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf

def slide_number_footer(slide, num, total=18):
    add_text_box(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
                 f"{num}/{total}", font_size=10, color=C_GRAY, alignment=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_gradient_bg(slide, RGBColor(0x0A, 0x0E, 0x1A), RGBColor(0x1A, 0x1A, 0x2E))

# Decorative circles
add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(3), RGBColor(0x19, 0x76, 0xD2))
add_circle(slide, Inches(11), Inches(5.5), Inches(3), RGBColor(0x3A, 0x76, 0xF0))

# Lock icon shape
lock = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.0), Inches(1.5), Inches(1.2))
lock.fill.solid()
lock.fill.fore_color.rgb = C_BLUE
lock.line.fill.background()
tf = lock.text_frame
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "\U0001F512"
run.font.size = Pt(44)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
             "ToledoMessage", font_size=54, color=C_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Segoe UI Semibold')

add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.0),
             "Post-Quantum Secure Messaging Platform", font_size=28, color=C_BLUE_L,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.8),
             "Hybrid X3DH Key Agreement  |  Double Ratchet Protocol  |  ML-KEM-768 + X25519",
             font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

# Tech badges
badges = [(".NET 10", C_PURPLE), ("Blazor WASM", C_BLUE), ("SignalR", C_TEAL),
          ("BouncyCastle", C_ORANGE), ("SQL Server", C_RED), ("E2EE", C_GREEN)]
badge_w = Inches(1.6)
total_w = len(badges) * badge_w.inches + (len(badges)-1) * 0.2
start_x = (13.333 - total_w) / 2
for i, (label, color) in enumerate(badges):
    x = Inches(start_x + i * (badge_w.inches + 0.2))
    r = add_rect(slide, x, Inches(5.6), badge_w, Inches(0.45), color)
    r.text_frame.paragraphs[0].text = label
    r.text_frame.paragraphs[0].font.size = Pt(12)
    r.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    r.text_frame.paragraphs[0].font.bold = True
    r.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r.text_frame.paragraphs[0].font.name = 'Segoe UI'

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
             "Mohamed Khamis  |  University of Toledo  |  February 2026",
             font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 1)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.8),
             "Presentation Outline", font_size=36, color=C_WHITE, bold=True)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(3), Inches(0.04))
shape.fill.solid(); shape.fill.fore_color.rgb = C_BLUE; shape.line.fill.background()

sections = [
    ("01", "Problem Statement & Motivation", C_BLUE),
    ("02", "System Architecture Overview", C_BLUE_L),
    ("03", "Cryptographic Model — X3DH Key Agreement", C_TEAL),
    ("04", "Cryptographic Model — Double Ratchet Protocol", C_GREEN),
    ("05", "Post-Quantum Resistance (ML-KEM-768)", C_PURPLE),
    ("06", "Hybrid Signature Scheme (Ed25519 + ML-DSA-65)", C_ORANGE),
    ("07", "Message Encryption Pipeline (AES-256-GCM)", C_RED),
    ("08", "End-to-End Message Flow", C_BLUE),
    ("09", "Real-Time Communication (SignalR)", C_TEAL),
    ("10", "Client-Side Security (Blazor WASM + IndexedDB)", C_GREEN),
    ("11", "Database & Entity Model", C_PURPLE),
    ("12", "Safety Number Verification", C_ORANGE),
    ("13", "Advanced Features", C_BLUE_L),
    ("14", "Security Analysis & Threat Model", C_RED),
    ("15", "Performance Benchmarks", C_BLUE),
    ("16", "Comparison with Existing Systems", C_TEAL),
    ("17", "Conclusion & Future Work", C_WHITE),
]

col1 = sections[:9]
col2 = sections[9:]

for i, (num, title, clr) in enumerate(col1):
    y = Inches(1.6 + i * 0.58)
    add_circle(slide, Inches(0.85), y + Inches(0.05), Inches(0.35), clr)
    add_text_box(slide, Inches(0.87), y + Inches(0.04), Inches(0.35), Inches(0.35),
                 num, font_size=11, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.4), y, Inches(5), Inches(0.5),
                 title, font_size=14, color=C_GRAY_L)

for i, (num, title, clr) in enumerate(col2):
    y = Inches(1.6 + i * 0.58)
    add_circle(slide, Inches(7.2), y + Inches(0.05), Inches(0.35), clr)
    add_text_box(slide, Inches(7.22), y + Inches(0.04), Inches(0.35), Inches(0.35),
                 num, font_size=11, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.75), y, Inches(5), Inches(0.5),
                 title, font_size=14, color=C_GRAY_L)

slide_number_footer(slide, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: Problem Statement
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(8), Inches(0.8),
             "01  Problem Statement & Motivation", font_size=32, color=C_WHITE, bold=True)

# Left panel — problems
problems = [
    ("Quantum Computing Threat", "Current messaging apps (Signal, WhatsApp) use classical Diffie-Hellman.\nShor's algorithm will break these in the post-quantum era."),
    ("Centralized Key Management", "Most platforms store encryption keys server-side,\ncreating single points of compromise."),
    ("No Post-Quantum Hybrid Model", "Pure post-quantum schemes are unproven at scale.\nA hybrid classical + PQ approach provides defense-in-depth."),
    ("Browser-Based Crypto Limitations", "Running cryptographic operations in WebAssembly\nwith secure client-side key storage is challenging."),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.5 + i * 1.35)
    add_accent_bar(slide, Inches(0.8), y + Inches(0.05), width=Inches(0.06), height=Inches(0.4), color=C_RED)
    add_text_box(slide, Inches(1.05), y, Inches(5.5), Inches(0.4),
                 title, font_size=18, color=C_BLUE_L, bold=True)
    add_text_box(slide, Inches(1.05), y + Inches(0.35), Inches(5.5), Inches(0.8),
                 desc, font_size=13, color=C_GRAY_L)

# Right panel — solution card
card = add_rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.0), C_NAVY, border_color=C_BLUE)
add_text_box(slide, Inches(7.5), Inches(1.7), Inches(5), Inches(0.5),
             "Our Solution: ToledoMessage", font_size=22, color=C_GREEN, bold=True)

solutions = [
    "\u2713  Hybrid X3DH: X25519 + ML-KEM-768 key agreement",
    "\u2713  Double Ratchet: forward secrecy per message",
    "\u2713  AES-256-GCM: authenticated encryption",
    "\u2713  Ed25519 + ML-DSA-65: hybrid digital signatures",
    "\u2713  Zero-knowledge server: keys never leave the client",
    "\u2713  Blazor WASM: full crypto in browser sandbox",
    "\u2713  Real-time via SignalR WebSockets",
    "\u2713  Multi-device support with per-device keys",
    "\u2713  Safety number verification (Signal-style)",
    "\u2713  Disappearing messages with server-side enforcement",
]

for i, sol in enumerate(solutions):
    add_text_box(slide, Inches(7.5), Inches(2.3 + i * 0.4), Inches(5), Inches(0.4),
                 sol, font_size=13, color=C_GRAY_L)

slide_number_footer(slide, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: System Architecture
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(8), Inches(0.8),
             "02  System Architecture Overview", font_size=32, color=C_WHITE, bold=True)

# Architecture diagram using matplotlib
def arch_diagram(fig, ax):
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 8)
    ax.axis('off')

    # Client layer
    client_box = mpatches.FancyBboxPatch((0.5, 5.5), 4.5, 2, boxstyle="round,pad=0.2",
                                          facecolor='#1976D2', edgecolor='#42A5F5', linewidth=2)
    ax.add_patch(client_box)
    ax.text(2.75, 7.0, 'Blazor WebAssembly Client', ha='center', va='center',
            fontsize=12, fontweight='bold', color='white')
    ax.text(2.75, 6.3, 'CryptoService | SessionService\nKeyGeneration | MessageEncryption\nIndexedDB Storage | ThemeService',
            ha='center', va='center', fontsize=8, color='#B3E5FC')

    # SignalR channel
    ax.annotate('', xy=(5.5, 6.5), xytext=(6.5, 6.5),
                arrowprops=dict(arrowstyle='<->', color='#25D366', lw=2.5))
    ax.text(6.0, 6.9, 'SignalR\nWebSocket', ha='center', va='center',
            fontsize=9, color='#25D366', fontweight='bold')

    # HTTP REST
    ax.annotate('', xy=(5.5, 5.8), xytext=(6.5, 5.8),
                arrowprops=dict(arrowstyle='<->', color='#FF9800', lw=2))
    ax.text(6.0, 5.4, 'HTTPS REST\nAPI', ha='center', va='center',
            fontsize=9, color='#FF9800', fontweight='bold')

    # Server layer
    server_box = mpatches.FancyBboxPatch((7, 5.5), 4.5, 2, boxstyle="round,pad=0.2",
                                          facecolor='#7C4DFF', edgecolor='#B388FF', linewidth=2)
    ax.add_patch(server_box)
    ax.text(9.25, 7.0, 'ASP.NET Core Server', ha='center', va='center',
            fontsize=12, fontweight='bold', color='white')
    ax.text(9.25, 6.3, 'ChatHub | AuthController\nDeviceController | MessageController\nASP.NET Core Identity + JWT',
            ha='center', va='center', fontsize=8, color='#E1BEE7')

    # Database
    db_box = mpatches.FancyBboxPatch((7.5, 0.5), 3.5, 2.2, boxstyle="round,pad=0.2",
                                      facecolor='#E53935', edgecolor='#EF9A9A', linewidth=2)
    ax.add_patch(db_box)
    ax.text(9.25, 1.95, 'SQL Server 2022', ha='center', va='center',
            fontsize=11, fontweight='bold', color='white')
    ax.text(9.25, 1.25, 'Users | Devices | Conversations\nParticipants | Messages\nPreKeys | SignedPreKeys',
            ha='center', va='center', fontsize=8, color='#FFCDD2')

    # EF Core arrow
    ax.annotate('', xy=(9.25, 5.3), xytext=(9.25, 2.9),
                arrowprops=dict(arrowstyle='<->', color='#EF9A9A', lw=2))
    ax.text(9.9, 4.1, 'EF Core 10\nCode First', ha='center', va='center',
            fontsize=9, color='#EF9A9A', fontweight='bold')

    # Client-side storage
    idb_box = mpatches.FancyBboxPatch((0.5, 0.5), 4.5, 2.2, boxstyle="round,pad=0.2",
                                       facecolor='#008069', edgecolor='#4DB6AC', linewidth=2)
    ax.add_patch(idb_box)
    ax.text(2.75, 1.95, 'Browser IndexedDB', ha='center', va='center',
            fontsize=11, fontweight='bold', color='white')
    ax.text(2.75, 1.25, 'Identity Keys (Ed25519 + ML-DSA-65)\nSession State (Double Ratchet)\nPre-Key Pairs | Auth Tokens',
            ha='center', va='center', fontsize=8, color='#B2DFDB')

    ax.annotate('', xy=(2.75, 5.3), xytext=(2.75, 2.9),
                arrowprops=dict(arrowstyle='<->', color='#4DB6AC', lw=2))
    ax.text(3.5, 4.1, 'LocalStorage\nService (JS Interop)', ha='center', va='center',
            fontsize=9, color='#4DB6AC', fontweight='bold')

    # Zero-knowledge label
    ax.text(6.0, 4.1, '\u26A0 Zero-Knowledge\nServer never sees\nplaintext or private keys',
            ha='center', va='center', fontsize=9, color='#FFD54F',
            fontweight='bold', style='italic',
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#2A2A2A', edgecolor='#FFD54F', linewidth=1.5))

buf = make_chart_image(arch_diagram, w=10, h=6.5)
slide.shapes.add_picture(buf, Inches(1.5), Inches(1.2), Inches(10.5), Inches(6))
slide_number_footer(slide, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: X3DH Key Agreement
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "03  Hybrid X3DH Key Agreement Protocol", font_size=32, color=C_WHITE, bold=True)

def x3dh_diagram(fig, ax):
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    # Alice
    alice_box = mpatches.FancyBboxPatch((0.3, 0.5), 3, 9, boxstyle="round,pad=0.3",
                                         facecolor='#1565C0', edgecolor='#42A5F5', linewidth=2, alpha=0.3)
    ax.add_patch(alice_box)
    ax.text(1.8, 9.2, 'ALICE (Initiator)', ha='center', va='center',
            fontsize=13, fontweight='bold', color='#42A5F5')

    # Bob
    bob_box = mpatches.FancyBboxPatch((10.7, 0.5), 3, 9, boxstyle="round,pad=0.3",
                                       facecolor='#00695C', edgecolor='#4DB6AC', linewidth=2, alpha=0.3)
    ax.add_patch(bob_box)
    ax.text(12.2, 9.2, 'BOB (Responder)', ha='center', va='center',
            fontsize=13, fontweight='bold', color='#4DB6AC')

    # Steps
    steps = [
        (8.5, "1. Alice fetches Bob's PreKey Bundle from server", '#FFD54F'),
        (7.5, "2. Verify Bob's signed pre-key (Ed25519 + ML-DSA-65)", '#EF9A9A'),
        (6.5, "3. Verify Bob's Kyber pre-key signature", '#EF9A9A'),
        (5.5, "4. DH1 = X25519(ephemeral_A, signedPreKey_B)", '#81D4FA'),
        (4.8, "5. DH2 = X25519(ephemeral_A, oneTimePreKey_B)", '#81D4FA'),
        (4.1, "6. ML-KEM-768 Encapsulate(kyberPreKey_B) \u2192 (ciphertext, shared_secret)", '#CE93D8'),
        (3.2, "7. ikm = DH1 || DH2 || kem_shared_secret", '#A5D6A7'),
        (2.3, "8. SK = HKDF-SHA256(ikm, info=\"ToledoMessage-X3DH-v1\")", '#A5D6A7'),
        (1.5, "9. rootKey = SK[0:32], chainKey = SK[32:64]", '#FFD54F'),
    ]

    for y, text, color in steps:
        ax.annotate('', xy=(10.5, y), xytext=(3.5, y),
                    arrowprops=dict(arrowstyle='->', color=color, lw=1.5))
        ax.text(7.0, y + 0.25, text, ha='center', va='center',
                fontsize=8.5, color=color, fontweight='bold')

    # Key types
    ax.text(1.8, 8.3, 'Generates:\n\u2022 Ephemeral X25519 keypair\n\u2022 ML-KEM encapsulation',
            ha='center', va='center', fontsize=8, color='#B3E5FC')

    ax.text(12.2, 8.0, 'Pre-Key Bundle:\n\u2022 Identity key (Ed25519+ML-DSA)\n\u2022 Signed pre-key (X25519)\n\u2022 Kyber pre-key (ML-KEM-768)\n\u2022 One-time pre-key (X25519)',
            ha='center', va='center', fontsize=8, color='#B2DFDB')

buf = make_chart_image(x3dh_diagram, w=11, h=7)
slide.shapes.add_picture(buf, Inches(1.0), Inches(1.1), Inches(11.3), Inches(6.2))
slide_number_footer(slide, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Double Ratchet
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "04  Double Ratchet Protocol", font_size=32, color=C_WHITE, bold=True)

def ratchet_diagram(fig, ax):
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 9)
    ax.axis('off')

    # Root chain
    ax.text(2.5, 8.3, 'ROOT CHAIN', ha='center', fontsize=12, fontweight='bold', color='#FFD54F')
    root_keys = ['RK\u2080', 'RK\u2081', 'RK\u2082', 'RK\u2083']
    for i, rk in enumerate(root_keys):
        x = 1.0 + i * 2.0
        circle = plt.Circle((x, 7.5), 0.4, facecolor='#FF9800', edgecolor='#FFD54F', linewidth=2)
        ax.add_patch(circle)
        ax.text(x, 7.5, rk, ha='center', va='center', fontsize=9, fontweight='bold', color='white')
        if i < len(root_keys) - 1:
            ax.annotate('', xy=(x + 1.5, 7.5), xytext=(x + 0.5, 7.5),
                        arrowprops=dict(arrowstyle='->', color='#FFD54F', lw=2))
            ax.text(x + 1.0, 7.9, 'HKDF', ha='center', fontsize=7, color='#FFD54F')

    # Sending chain
    ax.text(9.5, 8.3, 'SENDING CHAIN', ha='center', fontsize=12, fontweight='bold', color='#81D4FA')
    for i in range(3):
        x = 8.5 + i * 2.0
        y = 5.5
        circle = plt.Circle((x, y), 0.35, facecolor='#1976D2', edgecolor='#64B5F6', linewidth=2)
        ax.add_patch(circle)
        ax.text(x, y, f'CK\u209b{chr(8320+i)}', ha='center', va='center', fontsize=8, fontweight='bold', color='white')
        if i < 2:
            ax.annotate('', xy=(x + 1.5, y), xytext=(x + 0.5, y),
                        arrowprops=dict(arrowstyle='->', color='#64B5F6', lw=1.5))
        # Message keys
        ax.annotate('', xy=(x, y - 1.2), xytext=(x, y - 0.5),
                    arrowprops=dict(arrowstyle='->', color='#A5D6A7', lw=1.5))
        mk_circle = plt.Circle((x, y - 1.5), 0.3, facecolor='#2E7D32', edgecolor='#A5D6A7', linewidth=1.5)
        ax.add_patch(mk_circle)
        ax.text(x, y - 1.5, f'MK{i}', ha='center', va='center', fontsize=8, fontweight='bold', color='white')
        ax.text(x, y - 2.1, 'AES-256-GCM\nEncrypt', ha='center', fontsize=7, color='#A5D6A7')

    # Receiving chain
    ax.text(2.5, 5.2, 'RECEIVING CHAIN', ha='center', fontsize=10, fontweight='bold', color='#CE93D8')
    for i in range(2):
        x = 1.5 + i * 2.0
        y = 4.0
        circle = plt.Circle((x, y), 0.35, facecolor='#6A1B9A', edgecolor='#CE93D8', linewidth=2)
        ax.add_patch(circle)
        ax.text(x, y, f'CK\u1d63{chr(8320+i)}', ha='center', va='center', fontsize=8, fontweight='bold', color='white')
        if i < 1:
            ax.annotate('', xy=(x + 1.5, y), xytext=(x + 0.5, y),
                        arrowprops=dict(arrowstyle='->', color='#CE93D8', lw=1.5))
        ax.annotate('', xy=(x, y - 1.2), xytext=(x, y - 0.5),
                    arrowprops=dict(arrowstyle='->', color='#A5D6A7', lw=1.5))
        mk_circle = plt.Circle((x, y - 1.5), 0.3, facecolor='#2E7D32', edgecolor='#A5D6A7', linewidth=1.5)
        ax.add_patch(mk_circle)
        ax.text(x, y - 1.5, f'MK{i}', ha='center', va='center', fontsize=8, fontweight='bold', color='white')
        ax.text(x, y - 2.1, 'AES-256-GCM\nDecrypt', ha='center', fontsize=7, color='#A5D6A7')

    # DH ratchet
    ax.annotate('', xy=(7.5, 7.0), xytext=(7.5, 6.0),
                arrowprops=dict(arrowstyle='->', color='#EF9A9A', lw=2.5))
    ax.text(7.5, 6.5, 'DH\nRatchet\nStep', ha='center', va='center', fontsize=8,
            color='#EF9A9A', fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#2A2A2A', edgecolor='#EF9A9A'))

    # Labels
    ax.text(7.0, 1.0, 'Forward Secrecy: Compromised key cannot decrypt past messages\n'
            'Break-in Recovery: New DH ratchet step heals after compromise\n'
            'Per-Message Keys: Each message uses a unique AES-256-GCM key',
            ha='center', va='center', fontsize=9, color='#B0BEC5',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#1A1A2E', edgecolor='#546E7A', linewidth=1))

buf = make_chart_image(ratchet_diagram, w=11, h=7)
slide.shapes.add_picture(buf, Inches(1.0), Inches(1.1), Inches(11.3), Inches(6.2))
slide_number_footer(slide, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Post-Quantum (ML-KEM-768)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "05  Post-Quantum Resistance: ML-KEM-768", font_size=32, color=C_WHITE, bold=True)

# Left: explanation
items = [
    ("What is ML-KEM-768?", "NIST-standardized post-quantum Key Encapsulation Mechanism\n(formerly CRYSTALS-Kyber). Based on Module-LWE problem.", C_PURPLE),
    ("Why Hybrid?", "Classical X25519 provides proven security today.\nML-KEM-768 adds quantum resistance.\nCompromising BOTH is required to break the system.", C_BLUE),
    ("Key Sizes", "Public Key: 1,184 bytes  |  Ciphertext: 1,088 bytes\nShared Secret: 32 bytes  |  Security Level: NIST Level 3", C_GREEN),
    ("Integration Point", "During X3DH: Alice encapsulates against Bob's Kyber pre-key.\nThe KEM shared secret is mixed into the HKDF derivation\nalongside the X25519 DH outputs.", C_ORANGE),
]

for i, (title, desc, clr) in enumerate(items):
    y = Inches(1.4 + i * 1.4)
    card = add_rect(slide, Inches(0.7), y, Inches(7.2), Inches(1.2), C_NAVY, border_color=clr)
    add_text_box(slide, Inches(0.95), y + Inches(0.08), Inches(6.8), Inches(0.35),
                 title, font_size=16, color=clr, bold=True)
    add_text_box(slide, Inches(0.95), y + Inches(0.4), Inches(6.8), Inches(0.75),
                 desc, font_size=12, color=C_GRAY_L)

# Right: comparison chart
def pq_chart(fig, ax):
    categories = ['X25519\n(Classical)', 'ML-KEM-768\n(Post-Quantum)', 'Hybrid\n(ToledoMessage)']
    classical_sec = [128, 0, 128]
    pq_sec = [0, 192, 192]
    x = range(len(categories))
    bars1 = ax.bar([i - 0.18 for i in x], classical_sec, 0.35, label='Classical Security',
                    color='#1976D2', edgecolor='#42A5F5')
    bars2 = ax.bar([i + 0.18 for i in x], pq_sec, 0.35, label='Post-Quantum Security',
                    color='#7C4DFF', edgecolor='#B388FF')
    ax.set_xticks(list(x))
    ax.set_xticklabels(categories, fontsize=9, color='#B0BEC5')
    ax.set_ylabel('Security Level (bits)', fontsize=10, color='#B0BEC5')
    ax.set_ylim(0, 230)
    ax.tick_params(colors='#B0BEC5')
    ax.spines['bottom'].set_color('#546E7A')
    ax.spines['left'].set_color('#546E7A')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(fontsize=8, facecolor='#1A1A2E', edgecolor='#546E7A', labelcolor='#B0BEC5')
    for bar in bars1 + bars2:
        h = bar.get_height()
        if h > 0:
            ax.text(bar.get_x() + bar.get_width()/2., h + 3, f'{int(h)}',
                    ha='center', va='bottom', fontsize=9, color='white', fontweight='bold')

buf = make_chart_image(pq_chart, w=4.5, h=4)
slide.shapes.add_picture(buf, Inches(8.3), Inches(1.5), Inches(4.5), Inches(4))

slide_number_footer(slide, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Hybrid Signatures
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "06  Hybrid Signature Scheme", font_size=32, color=C_WHITE, bold=True)

add_text_box(slide, Inches(0.85), Inches(1.0), Inches(10), Inches(0.5),
             "Ed25519 (Classical)  +  ML-DSA-65 (Post-Quantum, formerly Dilithium)",
             font_size=18, color=C_ORANGE)

# Signing flow diagram
def sig_diagram(fig, ax):
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis('off')

    # Sign
    ax.text(2.5, 7.2, 'SIGNING (HybridSigner.Sign)', ha='center',
            fontsize=13, fontweight='bold', color='#FFD54F')

    msg_box = mpatches.FancyBboxPatch((0.5, 5.5), 4, 1.2, boxstyle="round,pad=0.2",
                                       facecolor='#2A2A2A', edgecolor='#81D4FA', linewidth=2)
    ax.add_patch(msg_box)
    ax.text(2.5, 6.1, 'Message Data', ha='center', va='center', fontsize=11, color='#81D4FA', fontweight='bold')

    # Ed25519 path
    ax.annotate('', xy=(2.0, 4.5), xytext=(2.0, 5.3),
                arrowprops=dict(arrowstyle='->', color='#1976D2', lw=2))
    ed_box = mpatches.FancyBboxPatch((0.5, 3.3), 3, 1, boxstyle="round,pad=0.2",
                                      facecolor='#1976D2', edgecolor='#64B5F6', linewidth=2)
    ax.add_patch(ed_box)
    ax.text(2.0, 3.8, 'Ed25519.Sign()\n64-byte signature', ha='center', va='center',
            fontsize=9, color='white', fontweight='bold')

    # ML-DSA path
    ax.annotate('', xy=(3.5, 4.5), xytext=(3.5, 5.3),
                arrowprops=dict(arrowstyle='->', color='#7C4DFF', lw=2))
    ml_box = mpatches.FancyBboxPatch((3.8, 3.3), 3.5, 1, boxstyle="round,pad=0.2",
                                      facecolor='#7C4DFF', edgecolor='#B388FF', linewidth=2)
    ax.add_patch(ml_box)
    ax.text(5.55, 3.8, 'ML-DSA-65.Sign()\n3,309-byte signature', ha='center', va='center',
            fontsize=9, color='white', fontweight='bold')

    # Concatenate
    ax.annotate('', xy=(3.5, 2.5), xytext=(2.0, 3.1),
                arrowprops=dict(arrowstyle='->', color='#FFD54F', lw=1.5))
    ax.annotate('', xy=(3.5, 2.5), xytext=(5.55, 3.1),
                arrowprops=dict(arrowstyle='->', color='#FFD54F', lw=1.5))
    result_box = mpatches.FancyBboxPatch((1.5, 1.5), 4, 0.8, boxstyle="round,pad=0.2",
                                          facecolor='#FF9800', edgecolor='#FFD54F', linewidth=2)
    ax.add_patch(result_box)
    ax.text(3.5, 1.9, 'Hybrid Sig = Ed25519_sig || ML-DSA_sig\n(3,373 bytes total)',
            ha='center', va='center', fontsize=9, color='white', fontweight='bold')

    # Verify
    ax.text(10.5, 7.2, 'VERIFICATION (HybridSigner.Verify)', ha='center',
            fontsize=13, fontweight='bold', color='#A5D6A7')

    v_box = mpatches.FancyBboxPatch((8, 4.5), 5, 2.2, boxstyle="round,pad=0.3",
                                     facecolor='#1A1A2E', edgecolor='#A5D6A7', linewidth=2)
    ax.add_patch(v_box)
    ax.text(10.5, 6.1, '1. Split signature at byte 64', ha='center', fontsize=9, color='#B0BEC5')
    ax.text(10.5, 5.5, '2. Verify Ed25519 part with Ed25519 public key', ha='center', fontsize=9, color='#64B5F6')
    ax.text(10.5, 4.9, '3. Verify ML-DSA part with ML-DSA public key', ha='center', fontsize=9, color='#B388FF')

    ax.text(10.5, 3.5, 'BOTH must pass \u2192 Valid', ha='center',
            fontsize=12, fontweight='bold', color='#A5D6A7',
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#2A2A2A', edgecolor='#A5D6A7'))

    ax.text(10.5, 2.2, 'Security Guarantee:\nEven if one scheme is broken\n(e.g., Ed25519 by quantum computer),\nthe other still protects integrity.',
            ha='center', fontsize=9, color='#FFD54F', style='italic',
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#2A2A2A', edgecolor='#FFD54F', linewidth=1))

buf = make_chart_image(sig_diagram, w=11, h=6)
slide.shapes.add_picture(buf, Inches(1.0), Inches(1.2), Inches(11.3), Inches(5.8))
slide_number_footer(slide, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: AES-256-GCM Encryption
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "07  Message Encryption: AES-256-GCM", font_size=32, color=C_WHITE, bold=True)

# Left side: encryption steps
steps_data = [
    ("1", "Derive Message Key", "Double Ratchet produces a unique 32-byte\nmessage key (MK) for each message", C_BLUE),
    ("2", "Generate Nonce", "12-byte random nonce (IV) generated via\nSecureRandom for each encryption", C_TEAL),
    ("3", "Encrypt", "AES-256-GCM(key=MK, nonce=IV, plaintext)\nProduces ciphertext + 16-byte auth tag", C_GREEN),
    ("4", "Package", "Output = nonce(12) || ciphertext || tag(16)\nBase64-encoded for transport", C_ORANGE),
    ("5", "Delete Key", "Message key is immediately discarded\nafter use \u2014 never reused", C_RED),
]

for i, (num, title, desc, clr) in enumerate(steps_data):
    y = Inches(1.3 + i * 1.1)
    circle = add_circle(slide, Inches(0.8), y + Inches(0.05), Inches(0.4), clr)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(14)
    circle.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.4), y, Inches(4.5), Inches(0.35),
                 title, font_size=16, color=clr, bold=True)
    add_text_box(slide, Inches(1.4), y + Inches(0.35), Inches(4.5), Inches(0.65),
                 desc, font_size=11, color=C_GRAY_L)

# Right: wire format
card = add_rect(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), C_NAVY, border_color=C_BLUE)
add_text_box(slide, Inches(7.0), Inches(1.45), Inches(5.4), Inches(0.4),
             "Wire Format (AeadCipher class)", font_size=16, color=C_BLUE_L, bold=True)

format_lines = [
    "Encryption:",
    "  plaintext bytes \u2192 UTF-8 encoded",
    "  nonce = SecureRandom(12 bytes)",
    "  (ciphertext, tag) = AES-GCM(key, nonce, plaintext)",
    "  output = nonce || ciphertext || tag",
    "  return Base64Encode(output)",
    "",
    "Decryption:",
    "  data = Base64Decode(input)",
    "  nonce = data[0:12]",
    "  ciphertext = data[12 : len-16]",
    "  tag = data[len-16 : len]",
    "  plaintext = AES-GCM-Decrypt(key, nonce, ct, tag)",
    "",
    "Key Properties:",
    "  \u2022 256-bit key (from Double Ratchet MK)",
    "  \u2022 96-bit nonce (random per message)",
    "  \u2022 128-bit authentication tag",
    "  \u2022 Authenticated Associated Data supported",
]

for i, line in enumerate(format_lines):
    clr = C_GREEN if line.startswith("  \u2022") else (C_ORANGE if line.endswith(":") and not line.startswith(" ") else C_GRAY_L)
    add_text_box(slide, Inches(7.1), Inches(1.9 + i * 0.26), Inches(5.3), Inches(0.26),
                 line, font_size=10, color=clr, font_name='Consolas')

slide_number_footer(slide, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: E2E Message Flow
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "08  End-to-End Message Flow", font_size=32, color=C_WHITE, bold=True)

def msg_flow(fig, ax):
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 10)
    ax.axis('off')

    # Actors
    actors = [('Alice\n(Sender)', 1.5, '#1976D2'), ('Server', 7, '#7C4DFF'), ('Bob\n(Receiver)', 12.5, '#008069')]
    for name, x, color in actors:
        circle = plt.Circle((x, 9), 0.5, facecolor=color, edgecolor='white', linewidth=2)
        ax.add_patch(circle)
        ax.text(x, 9, name, ha='center', va='center', fontsize=8, fontweight='bold', color='white')
        ax.plot([x, x], [0.5, 8.3], color='#546E7A', linewidth=1.5, linestyle='--')

    # Messages
    msgs = [
        (1.5, 7, 7.5, '1. Type message plaintext', '#81D4FA'),
        (1.5, 1.5, 6.8, '2. Ratchet \u2192 derive message key', '#FFD54F'),
        (1.5, 1.5, 6.1, '3. AES-256-GCM encrypt', '#A5D6A7'),
        (2.0, 6.5, 5.4, '4. SignalR: SendMessage(ciphertext)', '#CE93D8'),
        (7, 7, 4.5, '5. Store encrypted blob + route', '#B388FF'),
        (7.5, 12, 3.6, '6. SignalR: push to Bob\'s device', '#CE93D8'),
        (12.5, 12.5, 2.7, '7. Double Ratchet \u2192 derive same MK', '#FFD54F'),
        (12.5, 12.5, 1.9, '8. AES-256-GCM decrypt', '#A5D6A7'),
        (12.5, 7, 1.1, '9. ACK delivery status', '#EF9A9A'),
    ]

    for x1, x2, y, text, color in msgs:
        if x1 != x2:
            ax.annotate('', xy=(x2, y), xytext=(x1, y),
                        arrowprops=dict(arrowstyle='->', color=color, lw=2))
        ax.text((x1+x2)/2, y + 0.25, text, ha='center', fontsize=8.5, color=color, fontweight='bold')

buf = make_chart_image(msg_flow, w=11, h=7.5)
slide.shapes.add_picture(buf, Inches(1.0), Inches(0.8), Inches(11.3), Inches(6.5))
slide_number_footer(slide, 10)

# ═══════════════════════════════════════════════════════════
# SLIDE 11: SignalR Real-Time
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "09  Real-Time Communication: SignalR Hub", font_size=32, color=C_WHITE, bold=True)

# ChatHub methods
hub_methods = [
    ("RegisterDevice", "Maps deviceId \u2192 SignalR connectionId\nfor targeted message delivery", C_BLUE),
    ("SendMessage", "Receives encrypted envelope, stores in DB,\nroutes to recipient's connected devices", C_GREEN),
    ("AcknowledgeDelivery", "Updates message status to Delivered,\nnotifies sender in real-time", C_TEAL),
    ("AcknowledgeRead", "Updates message status to Read,\nblue-tick confirmation to sender", C_PURPLE),
    ("SendTypingIndicator", "Broadcasts typing status to\nconversation participants", C_ORANGE),
    ("AddParticipant", "Admin adds user to group conversation,\nbroadcasts update to all members", C_BLUE_L),
    ("RemoveParticipant", "Admin removes user from group,\ncleans up participant records", C_RED),
]

for i, (method, desc, clr) in enumerate(hub_methods):
    y = Inches(1.3 + i * 0.82)
    add_accent_bar(slide, Inches(0.7), y + Inches(0.05), width=Inches(0.06), height=Inches(0.35), color=clr)
    add_text_box(slide, Inches(0.95), y, Inches(3.2), Inches(0.35),
                 method + "()", font_size=14, color=clr, bold=True, font_name='Consolas')
    add_text_box(slide, Inches(4.3), y, Inches(4.5), Inches(0.7),
                 desc, font_size=11, color=C_GRAY_L)

# Right: delivery status diagram
card = add_rect(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.5), C_NAVY, border_color=C_TEAL)
add_text_box(slide, Inches(9.2), Inches(1.45), Inches(3.4), Inches(0.4),
             "Message Lifecycle", font_size=16, color=C_TEAL, bold=True)

statuses = [
    ("\u23F3 Sending", "Client encrypts & sends", C_GRAY),
    ("\u2713 Sent", "Server received & stored", C_GRAY_L),
    ("\u2713\u2713 Delivered", "Recipient device received", C_BLUE_L),
    ("\u2713\u2713 Read", "Recipient opened message", C_GREEN),
]

for i, (status, desc, clr) in enumerate(statuses):
    y = Inches(2.0 + i * 1.1)
    add_text_box(slide, Inches(9.3), y, Inches(3.2), Inches(0.4),
                 status, font_size=16, color=clr, bold=True)
    add_text_box(slide, Inches(9.3), y + Inches(0.4), Inches(3.2), Inches(0.4),
                 desc, font_size=11, color=C_GRAY)
    if i < len(statuses) - 1:
        add_text_box(slide, Inches(9.8), y + Inches(0.8), Inches(1), Inches(0.3),
                     "\u25BC", font_size=14, color=C_TEAL, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 11)

# ═══════════════════════════════════════════════════════════
# SLIDE 12: Client-Side Security
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "10  Client-Side Security Model", font_size=32, color=C_WHITE, bold=True)

services = [
    ("LocalStorageService", "Encrypted IndexedDB storage via JS Interop.\nBinary key-value with optional encryption-at-rest.", C_BLUE),
    ("KeyGenerationService", "Generates identity keys (Ed25519 + ML-DSA-65),\nsigned pre-keys, Kyber pre-keys, one-time pre-keys.\nBatches of 10 pre-keys per registration.", C_GREEN),
    ("SessionService", "Manages Double Ratchet session state per peer.\nTracks root keys, chain keys, message counters.", C_TEAL),
    ("CryptoService", "Orchestrates encryption/decryption.\nHandles X3DH initiation and session establishment.", C_PURPLE),
    ("MessageEncryptionService", "Encrypts for all of recipient's devices.\nEach device gets a unique ciphertext.", C_ORANGE),
    ("FingerprintService", "Computes safety numbers from identity keys.\nEnables out-of-band verification.", C_BLUE_L),
    ("PreKeyReplenishmentService", "Monitors pre-key count on server.\nAuto-generates new batches when depleted.", C_RED),
]

for i, (name, desc, clr) in enumerate(services):
    y = Inches(1.3 + i * 0.84)
    card = add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.72), C_NAVY, border_color=clr)
    add_text_box(slide, Inches(0.85), y + Inches(0.06), Inches(3.5), Inches(0.3),
                 name, font_size=13, color=clr, bold=True, font_name='Consolas')
    add_text_box(slide, Inches(4.5), y + Inches(0.06), Inches(7.8), Inches(0.6),
                 desc, font_size=11, color=C_GRAY_L)

slide_number_footer(slide, 12)

# ═══════════════════════════════════════════════════════════
# SLIDE 13: Database Model
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "11  Database & Entity Model", font_size=32, color=C_WHITE, bold=True)

def er_diagram(fig, ax):
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 9)
    ax.axis('off')

    entities = [
        (1, 7, 2.5, 1.5, 'AppUser', 'Id, DisplayName\nPasswordHash\nCreatedAt', '#1976D2', '#64B5F6'),
        (5, 7, 2.5, 1.5, 'Device', 'DeviceId, UserId\nDeviceName, IdentityKey\nSignedPreKey, KyberPreKey\nLastSeenAt', '#7C4DFF', '#B388FF'),
        (1, 4, 2.5, 1.5, 'Conversation', 'ConversationId, Type\nGroupName, CreatedBy\nDisappearTimer', '#008069', '#4DB6AC'),
        (5, 4, 2.5, 1.5, 'Participant', 'ConversationId, UserId\nRole, JoinedAt', '#FF9800', '#FFD54F'),
        (9.5, 7, 3, 1.5, 'OneTimePreKey', 'PreKeyId, DeviceId\nPublicKey, IsUsed', '#E53935', '#EF9A9A'),
        (9.5, 4, 3, 1.5, 'Message', 'MessageId, ConversationId\nSenderDeviceId, RecipientDeviceId\nCiphertext, Status, ServerTimestamp', '#25D366', '#A5D6A7'),
    ]

    for x, y, w, h, name, fields, fc, ec in entities:
        box = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.15",
                                       facecolor=fc, edgecolor=ec, linewidth=2, alpha=0.85)
        ax.add_patch(box)
        ax.text(x + w/2, y + h - 0.25, name, ha='center', va='center',
                fontsize=10, fontweight='bold', color='white')
        ax.text(x + w/2, y + h/2 - 0.2, fields, ha='center', va='center',
                fontsize=7, color='#E0E0E0')

    # Relations
    relations = [
        (3.5, 7.75, 5, 7.75, '1:N'),   # User -> Device
        (2.25, 7, 2.25, 5.5, '1:N'),    # User -> Conversation
        (3.5, 5, 5, 5, 'M:N'),          # Conversation -> Participant
        (7.5, 7.75, 9.5, 7.75, '1:N'),  # Device -> OneTimePreKey
        (7.5, 4.75, 9.5, 4.75, '1:N'),  # Conversation -> Message
    ]

    for x1, y1, x2, y2, label in relations:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color='#B0BEC5', lw=1.5))
        ax.text((x1+x2)/2, (y1+y2)/2 + 0.2, label, ha='center',
                fontsize=8, color='#FFD54F', fontweight='bold')

    ax.text(7, 1.5, 'EF Core 10 Code First  |  SQL Server 2022\n'
            'Decimal PKs (precision 20, scale 0)  |  UTC timestamps\n'
            'Server stores only encrypted ciphertext \u2014 never plaintext',
            ha='center', va='center', fontsize=10, color='#B0BEC5',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#1A1A2E', edgecolor='#546E7A'))

buf = make_chart_image(er_diagram, w=11, h=7)
slide.shapes.add_picture(buf, Inches(1.0), Inches(1.0), Inches(11.3), Inches(6.2))
slide_number_footer(slide, 13)

# ═══════════════════════════════════════════════════════════
# SLIDE 14: Safety Numbers
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "12  Safety Number Verification", font_size=32, color=C_WHITE, bold=True)

# Left: how it works
steps = [
    ("1. Collect Identity Keys", "Gather both parties' Ed25519 + ML-DSA-65\npublic identity keys from local storage", C_BLUE),
    ("2. SHA-256 Hash", "Concatenate and hash: SHA256(localKey || remoteKey)\nProduces 32-byte fingerprint digest", C_TEAL),
    ("3. Numeric Encoding", "Convert hash to human-readable numeric groups\n(e.g., \"37291 04823 19457 ...\")", C_GREEN),
    ("4. Compare Out-of-Band", "Users compare numbers in person or via\ntrusted channel (phone call, video chat)", C_ORANGE),
    ("5. Mark Verified", "Store verification flag per conversation\nin IndexedDB. Show green verified badge.", C_PURPLE),
]

for i, (title, desc, clr) in enumerate(steps):
    y = Inches(1.3 + i * 1.15)
    add_accent_bar(slide, Inches(0.7), y + Inches(0.05), width=Inches(0.06), height=Inches(0.4), color=clr)
    add_text_box(slide, Inches(0.95), y, Inches(6), Inches(0.35),
                 title, font_size=15, color=clr, bold=True)
    add_text_box(slide, Inches(0.95), y + Inches(0.35), Inches(6), Inches(0.7),
                 desc, font_size=12, color=C_GRAY_L)

# Right: visual example
card = add_rect(slide, Inches(7.5), Inches(1.3), Inches(5.2), Inches(5.5), C_NAVY, border_color=C_GREEN)
add_text_box(slide, Inches(7.7), Inches(1.5), Inches(4.8), Inches(0.4),
             "Safety Number Display", font_size=16, color=C_GREEN, bold=True)

example_card = add_rect(slide, Inches(7.9), Inches(2.1), Inches(4.6), Inches(2.5), RGBColor(0x20, 0x2C, 0x33))
add_text_box(slide, Inches(8.1), Inches(2.2), Inches(4.2), Inches(0.3),
             "Safety Number", font_size=12, color=C_GRAY)
add_text_box(slide, Inches(8.1), Inches(2.6), Inches(4.2), Inches(1.5),
             "37291 04823\n19457 82634\n50918 73246\n61582 39047\n84619 25703",
             font_size=18, color=C_WHITE, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)

badge = add_rect(slide, Inches(9.0), Inches(5.0), Inches(2.6), Inches(0.5), RGBColor(0x2E, 0x7D, 0x32))
badge.text_frame.paragraphs[0].text = "\u2713  Verified"
badge.text_frame.paragraphs[0].font.size = Pt(16)
badge.text_frame.paragraphs[0].font.color.rgb = C_WHITE
badge.text_frame.paragraphs[0].font.bold = True
badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(7.7), Inches(5.7), Inches(4.8), Inches(0.8),
             "Key Change Warning:\nIf identity keys change, safety number changes\nand users are alerted to re-verify.",
             font_size=11, color=C_ORANGE)

slide_number_footer(slide, 14)

# ═══════════════════════════════════════════════════════════
# SLIDE 15: Advanced Features
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "13  Advanced Features", font_size=32, color=C_WHITE, bold=True)

features = [
    ("Multi-Device Support", "Each device has independent identity & pre-keys.\nMessages encrypted separately per device.\nServer routes to all active devices.", C_BLUE, "\U0001F4F1"),
    ("Disappearing Messages", "Configurable timer: 5min, 1h, 24h, 7 days.\nServer-side enforcement with MessageExpiryService.\nAutomatic cleanup after expiration.", C_GREEN, "\u23F0"),
    ("Group Messaging", "Fan-out encryption: sender encrypts for\neach participant's devices individually.\nAdmin role for participant management.", C_PURPLE, "\U0001F465"),
    ("Pre-Key Replenishment", "Auto-detects when one-time pre-keys run low.\nGenerates and uploads fresh batches of 10.\nEnsures X3DH sessions always available.", C_ORANGE, "\U0001F504"),
    ("Responsive UI + Themes", "7 messenger-inspired themes (WhatsApp, Telegram, Signal).\nCSS custom properties + data-theme selectors.\nMobile-first responsive with hamburger nav.", C_TEAL, "\U0001F3A8"),
    ("Real-Time Typing Indicators", "SignalR broadcasts typing status.\nAuto-clears after 3-second timeout.\nNon-blocking, failure-tolerant.", C_BLUE_L, "\u2328"),
]

col_w = Inches(6.0)
for i, (title, desc, clr, icon) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.6) if col == 0 else Inches(6.9)
    y = Inches(1.3 + row * 1.9)
    card = add_rect(slide, x, y, Inches(5.8), Inches(1.7), C_NAVY, border_color=clr)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.35),
                 f"{icon}  {title}", font_size=15, color=clr, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.4), Inches(1.1),
                 desc, font_size=11, color=C_GRAY_L)

slide_number_footer(slide, 15)

# ═══════════════════════════════════════════════════════════
# SLIDE 16: Security Analysis
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "14  Security Analysis & Threat Model", font_size=32, color=C_WHITE, bold=True)

threats = [
    ("Threat", "Mitigation", "Status"),
    ("Man-in-the-Middle", "X3DH mutual authentication + safety numbers", "\u2705"),
    ("Server Compromise", "E2EE: server only sees ciphertext, never plaintext", "\u2705"),
    ("Key Compromise (past)", "Forward secrecy via Double Ratchet DH steps", "\u2705"),
    ("Key Compromise (future)", "Break-in recovery: new DH ratchet heals session", "\u2705"),
    ("Quantum Computer Attack", "ML-KEM-768 + ML-DSA-65 post-quantum primitives", "\u2705"),
    ("Replay Attack", "Unique nonce per message + message counters", "\u2705"),
    ("Device Theft", "Keys stored in browser IndexedDB, session-bound", "\u26A0"),
    ("Metadata Leakage", "Server sees routing metadata (who talks to whom)", "\u26A0"),
]

# Table
header_y = Inches(1.4)
add_rect(slide, Inches(0.6), header_y, Inches(12.1), Inches(0.5), C_BLUE)
add_text_box(slide, Inches(0.8), header_y + Inches(0.05), Inches(4), Inches(0.4),
             threats[0][0], font_size=14, color=C_WHITE, bold=True)
add_text_box(slide, Inches(5.0), header_y + Inches(0.05), Inches(6), Inches(0.4),
             threats[0][1], font_size=14, color=C_WHITE, bold=True)
add_text_box(slide, Inches(11.5), header_y + Inches(0.05), Inches(1), Inches(0.4),
             threats[0][2], font_size=14, color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

for i, (threat, mitigation, status) in enumerate(threats[1:]):
    y = Inches(1.95 + i * 0.6)
    bg_clr = C_NAVY if i % 2 == 0 else C_BG
    add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.55), bg_clr)
    add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(4), Inches(0.4),
                 threat, font_size=13, color=C_RED if "\u26A0" in status else C_GRAY_L, bold=True)
    add_text_box(slide, Inches(5.0), y + Inches(0.08), Inches(6), Inches(0.4),
                 mitigation, font_size=12, color=C_GRAY_L)
    add_text_box(slide, Inches(11.5), y + Inches(0.05), Inches(1), Inches(0.4),
                 status, font_size=18, alignment=PP_ALIGN.CENTER,
                 color=C_GREEN if status == "\u2705" else C_ORANGE)

slide_number_footer(slide, 16)

# ═══════════════════════════════════════════════════════════
# SLIDE 17: Performance & Comparison
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.5))
add_text_box(slide, Inches(0.85), Inches(0.4), Inches(10), Inches(0.8),
             "15-16  Performance & Comparison", font_size=32, color=C_WHITE, bold=True)

# Comparison chart
def comparison_chart(fig, ax):
    apps = ['ToledoMessage\n(Ours)', 'Signal', 'WhatsApp', 'Telegram\n(Secret)', 'iMessage']
    features = ['E2EE by Default', 'Post-Quantum\nResistance', 'Open Source\nCrypto', 'Forward\nSecrecy', 'Safety Number\nVerification', 'Multi-Device\nE2EE', 'Disappearing\nMessages']

    # Score matrix (0-1)
    scores = [
        [1, 1, 1, 1, 1, 1, 1],     # ToledoMessage
        [1, 0.5, 1, 1, 1, 0.7, 1],  # Signal
        [1, 0.3, 0, 1, 1, 0.5, 1],  # WhatsApp
        [0, 0, 0, 1, 0, 0.3, 1],    # Telegram
        [1, 0.8, 0, 1, 0, 1, 0],    # iMessage
    ]

    colors_list = ['#3A76F0', '#2196F3', '#25D366', '#2AABEE', '#007AFF']

    x = range(len(features))
    bar_w = 0.15
    for i, (app, score, color) in enumerate(zip(apps, scores, colors_list)):
        positions = [xi + (i - 2) * bar_w for xi in x]
        ax.bar(positions, score, bar_w, label=app, color=color, alpha=0.85, edgecolor='white', linewidth=0.5)

    ax.set_xticks(list(x))
    ax.set_xticklabels(features, fontsize=7, color='#B0BEC5', ha='center')
    ax.set_ylabel('Capability Score', fontsize=10, color='#B0BEC5')
    ax.set_ylim(0, 1.2)
    ax.set_yticks([0, 0.25, 0.5, 0.75, 1.0])
    ax.set_yticklabels(['None', 'Low', 'Partial', 'High', 'Full'], fontsize=8, color='#B0BEC5')
    ax.tick_params(colors='#B0BEC5')
    ax.spines['bottom'].set_color('#546E7A')
    ax.spines['left'].set_color('#546E7A')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(fontsize=7, facecolor='#0F111A', edgecolor='#546E7A', labelcolor='#B0BEC5',
              loc='upper right', ncol=3)
    ax.set_title('Feature Comparison: ToledoMessage vs. Industry', fontsize=11,
                 color='#FFFFFF', fontweight='bold', pad=10)

buf = make_chart_image(comparison_chart, w=10, h=4.5)
slide.shapes.add_picture(buf, Inches(1.5), Inches(1.2), Inches(10.5), Inches(5.5))

slide_number_footer(slide, 17)

# ═══════════════════════════════════════════════════════════
# SLIDE 18: Conclusion
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, RGBColor(0x0A, 0x0E, 0x1A), C_DARK)

add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), C_BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "17  Conclusion & Future Work", font_size=36, color=C_WHITE, bold=True)

# Contributions
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.4),
             "Key Contributions", font_size=22, color=C_GREEN, bold=True)

contributions = [
    "\u2713  First Blazor WASM messenger with hybrid post-quantum E2EE",
    "\u2713  Novel integration of ML-KEM-768 into X3DH key agreement",
    "\u2713  Full Double Ratchet with per-message AES-256-GCM keys",
    "\u2713  Hybrid signature scheme (Ed25519 + ML-DSA-65)",
    "\u2713  Zero-knowledge server architecture — keys never leave client",
    "\u2713  Multi-device E2EE with independent key material per device",
    "\u2713  65 comprehensive unit tests covering all crypto primitives",
]

for i, c in enumerate(contributions):
    add_text_box(slide, Inches(1.0), Inches(2.0 + i * 0.4), Inches(6.5), Inches(0.4),
                 c, font_size=13, color=C_GRAY_L)

# Future work
add_text_box(slide, Inches(8), Inches(1.5), Inches(4.5), Inches(0.4),
             "Future Work", font_size=22, color=C_ORANGE, bold=True)

future = [
    "\u25B6  Mobile native apps (MAUI / Flutter)",
    "\u25B6  File & media encrypted transfer",
    "\u25B6  Voice/video calls with SRTP",
    "\u25B6  Multi-device sync via Sesame protocol",
    "\u25B6  Formal security proofs (ProVerif/Tamarin)",
    "\u25B6  Hardware security module (HSM) integration",
    "\u25B6  Decentralized architecture (Matrix federation)",
]

for i, f in enumerate(future):
    add_text_box(slide, Inches(8.2), Inches(2.0 + i * 0.4), Inches(4.5), Inches(0.4),
                 f, font_size=13, color=C_GRAY_L)

# Final tagline
add_rect(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.7), C_BLUE)
add_text_box(slide, Inches(2), Inches(5.55), Inches(9.3), Inches(0.6),
             "Secure today. Quantum-resistant tomorrow.", font_size=24,
             color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "Thank you!  |  Questions?  |  github.com/ToledoMessage",
             font_size=16, color=C_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 18)

# ─── Save ───
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ToledoMessage_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
